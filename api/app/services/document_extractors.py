"""
文档格式提取器

支持 PDF、TXT、DOCX、PPTX 等多种格式的文本提取
"""
from typing import Optional, Dict, Any
from pathlib import Path
import logging

logger = logging.getLogger(__name__)


class DocumentExtractor:
    """文档提取器基类"""

    def extract_text(self, file_path: str) -> str:
        """
        提取文档文本内容

        Args:
            file_path: 文件路径

        Returns:
            str: 提取的文本内容
        """
        raise NotImplementedError("Subclasses must implement extract_text")


class PDFExtractor(DocumentExtractor):
    """PDF 文档提取器"""

    def extract_text(self, file_path: str) -> str:
        """从 PDF 提取文本"""
        try:
            import fitz  # PyMuPDF
            text_parts = []

            with fitz.open(file_path) as doc:
                for page in doc:
                    text = page.get_text()
                    if text.strip():
                        text_parts.append(text)

            full_text = "\n\n".join(text_parts)
            logger.info(f"✅ PDF 文本提取成功: {len(full_text)} 字符")
            return full_text

        except Exception as e:
            logger.error(f"❌ PDF 提取失败: {e}")
            raise


class TXTExtractor(DocumentExtractor):
    """TXT 文本提取器"""

    def extract_text(self, file_path: str) -> str:
        """从 TXT 文件提取文本"""
        try:
            # 尝试多种编码
            for encoding in ['utf-8', 'gbk', 'gb2312', 'latin-1']:
                try:
                    with open(file_path, 'r', encoding=encoding) as f:
                        content = f.read()
                    logger.info(f"✅ TXT 文本提取成功 ({encoding}): {len(content)} 字符")
                    return content
                except UnicodeDecodeError:
                    continue

            # 如果所有编码都失败，使用 utf-8 并忽略错误
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                content = f.read()
            logger.warning(f"⚠️ TXT 文本提取成功(忽略错误): {len(content)} 字符")
            return content

        except Exception as e:
            logger.error(f"❌ TXT 提取失败: {e}")
            raise


class DocxExtractor(DocumentExtractor):
    """Word 文档提取器 (.docx)"""

    def extract_text(self, file_path: str) -> str:
        """从 Word 文档提取文本"""
        try:
            from docx import Document

            doc = Document(file_path)
            text_parts = []

            # 提取段落文本
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text_parts.append(paragraph.text)

            # 提取表格文本
            for table in doc.tables:
                for row in table.rows:
                    row_text = " | ".join([cell.text for cell in row.cells])
                    if row_text.strip():
                        text_parts.append(row_text)

            full_text = "\n\n".join(text_parts)
            logger.info(f"✅ Word 文档文本提取成功: {len(full_text)} 字符")
            return full_text

        except Exception as e:
            logger.error(f"❌ Word 文档提取失败: {e}")
            raise


class PptxExtractor(DocumentExtractor):
    """PowerPoint 提取器 (.pptx)"""

    def extract_text(self, file_path: str) -> str:
        """从 PowerPoint 提取文本"""
        try:
            from pptx import Presentation

            prs = Presentation(file_path)
            text_parts = []

            # 提取幻灯片中的文本
            for slide_num, slide in enumerate(prs.slides, start=1):
                slide_texts = []

                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_texts.append(shape.text)

                if slide_texts:
                    text_parts.append(f"[幻灯片 {slide_num}]")
                    text_parts.extend(slide_texts)
                    text_parts.append("")  # 幻灯片之间空行

            full_text = "\n".join(text_parts)
            logger.info(f"✅ PowerPoint 文本提取成功: {len(full_text)} 字符")
            return full_text

        except Exception as e:
            logger.error(f"❌ PowerPoint 提取失败: {e}")
            raise


class DocumentExtractorFactory:
    """文档提取器工厂类"""

    _extractors = {
        'pdf': PDFExtractor,
        'txt': TXTExtractor,
        'text': TXTExtractor,
        'docx': DocxExtractor,
        'docm': DocxExtractor,  # .docm 也是 docx 格式
        'pptx': PptxExtractor,
        'ppt': PptxExtractor,   # .ppt 实际上需要不同的处理，这里暂用 pptx
    }

    @classmethod
    def get_extractor(cls, file_type: str) -> Optional[DocumentExtractor]:
        """
        根据文件类型获取对应的提取器

        Args:
            file_type: 文件类型（不含点号，如 'pdf', 'docx'）

        Returns:
            DocumentExtractor: 提取器实例，如果不支持的类型则返回 None
        """
        file_type = file_type.lower().strip()

        if file_type in cls._extractors:
            return cls._extractors[file_type]()

        return None

    @classmethod
    def supported_formats(cls) -> list:
        """获取支持的文件格式列表"""
        return list(cls._extractors.keys())

    @classmethod
    def is_supported(cls, file_type: str) -> bool:
        """
        检查文件类型是否支持

        Args:
            file_type: 文件类型（可以包含点号）

        Returns:
            bool: 是否支持
        """
        # 去掉点号
        file_type = file_type.lower().replace('.', '').strip()
        return file_type in cls._extractors


def extract_text_from_file(file_path: str) -> str:
    """
    自动识别文件类型并提取文本

    Args:
        file_path: 文件路径

    Returns:
        str: 提取的文本内容

    Raises:
        ValueError: 不支持的文件类型
        Exception: 文本提取失败
    """
    # 从文件扩展名获取类型
    path_obj = Path(file_path)
    file_ext = path_obj.suffix.lstrip('.')

    # 获取提取器
    extractor = DocumentExtractorFactory.get_extractor(file_ext)

    if not extractor:
        supported = ", ".join(DocumentExtractorFactory.supported_formats())
        raise ValueError(
            f"不支持的文件格式: .{file_ext}。支持的格式: {supported}"
        )

    # 提取文本
    return extractor.extract_text(file_path)
